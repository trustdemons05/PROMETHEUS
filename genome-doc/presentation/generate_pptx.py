"""
Generate Genome-Doc Academic Project Presentation (17 Slides)
Run: python presentation/generate_pptx.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree
from pathlib import Path

# ── Colors ──
NAVY      = RGBColor(0x1B, 0x2A, 0x4A)
NAVY_DK   = RGBColor(0x0F, 0x17, 0x2A)
TEAL      = RGBColor(0x0D, 0x94, 0x88)
INDIGO    = RGBColor(0x4F, 0x46, 0xE5)
ROSE      = RGBColor(0xBE, 0x12, 0x3C)
SL8       = RGBColor(0x1E, 0x29, 0x3B)
SL6       = RGBColor(0x47, 0x55, 0x69)
SL4       = RGBColor(0x94, 0xA3, 0xB8)
SL2       = RGBColor(0xE2, 0xE8, 0xF0)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GREEN7    = RGBColor(0x15, 0x80, 0x3D)
GREEN1    = RGBColor(0xDC, 0xFC, 0xE7)
AMBER6    = RGBColor(0xD9, 0x77, 0x06)
AMBER1    = RGBColor(0xFE, 0xF3, 0xC7)

SW = Inches(13.333)
SH = Inches(7.5)
DIR = Path(__file__).resolve().parent
ASSETS = DIR / "assets"
TOTAL = 17

# ── Transitions ──
TR = {
    "fade":   ("fade",   {}),
    "push":   ("push",   {"dir": "l"}),
    "wipe":   ("wipe",   {"dir": "d"}),
    "cover":  ("cover",  {"dir": "l"}),
    "split":  ("split",  {"orient": "horz", "dir": "out"}),
    "reveal": ("reveal", {"dir": "l"}),
}

def add_tr(slide, name="fade"):
    t, a = TR.get(name, TR["fade"])
    el = etree.SubElement(slide._element, qn("p:transition"))
    el.set("spd", "med"); el.set("advClick", "1")
    ch = etree.SubElement(el, qn(f"p:{t}"))
    for k, v in a.items(): ch.set(k, v)

# ── Helpers ──
def bg(s, c):
    s.background.fill.solid(); s.background.fill.fore_color.rgb = c

def bar(s, c=NAVY, y=0):
    b = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, y, SW, Inches(0.06))
    b.fill.solid(); b.fill.fore_color.rgb = c; b.line.fill.background()

def snum(s, n):
    t = s.shapes.add_textbox(Inches(12.3), Inches(7.0), Inches(0.8), Inches(0.3))
    p = t.text_frame.paragraphs[0]; p.text = f"{n} / {TOTAL}"
    p.font.size = Pt(9); p.font.color.rgb = SL4; p.font.name = "Calibri"
    p.alignment = PP_ALIGN.RIGHT

def tx(s, l, t, w, h, text, sz=16, c=SL8, b=False, al=PP_ALIGN.LEFT, fn="Calibri"):
    tb = s.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(sz); p.font.color.rgb = c; p.font.bold = b
    p.font.name = fn; p.alignment = al
    return tb

def ml(s, l, t, w, h, lines, sz=12, c=SL6, sp=0.45):
    """lines: list of str or (text, color, bold, size)"""
    tb = s.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, ld in enumerate(lines):
        if isinstance(ld, str): txt, cl, bd, fs = ld, c, False, sz
        else:
            txt = ld[0]; cl = ld[1] if len(ld)>1 else c
            bd = ld[2] if len(ld)>2 else False; fs = ld[3] if len(ld)>3 else sz
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = txt; p.font.size = Pt(fs); p.font.color.rgb = cl
        p.font.bold = bd; p.font.name = "Calibri"; p.space_after = Pt(fs * sp)
    return tb

def tag(s, l, t, text, fg, bgc, w=Inches(1.4)):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, Inches(0.3))
    sh.fill.solid(); sh.fill.fore_color.rgb = bgc
    sh.line.color.rgb = fg; sh.line.width = Pt(0.8)
    p = sh.text_frame.paragraphs[0]; p.text = text
    p.font.size = Pt(8); p.font.color.rgb = fg; p.font.bold = True
    p.font.name = "Calibri"; p.alignment = PP_ALIGN.CENTER
    sh.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

def crd(s, l, t, w, h, bc=SL2):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = WHITE
    sh.line.color.rgb = bc; sh.line.width = Pt(1.2); return sh

def topbar(s, x, y, w, col):
    b = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Inches(0.05))
    b.fill.solid(); b.fill.fore_color.rgb = col; b.line.fill.background()

def header(s, label, title, label_color=TEAL, n=1, tr="fade"):
    bg(s, WHITE); bar(s, NAVY); add_tr(s, tr)
    tx(s, Inches(0.8), Inches(0.4), Inches(3), Inches(0.3),
       label, sz=10, c=label_color, b=True)
    tx(s, Inches(0.8), Inches(0.9), Inches(11), Inches(0.6),
       title, sz=32, c=NAVY, b=True)
    snum(s, n)

def bullets(s, l, t, w, items, sz=12, c=SL6, gap=0.4):
    tb = s.shapes.add_textbox(l, t, w, Inches(len(items) * gap + 0.5))
    tf = tb.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"•  {item}"
        p.font.size = Pt(sz); p.font.color.rgb = c; p.font.name = "Calibri"
        p.space_after = Pt(sz * 0.35)

def add_img(s, name, l, t, w, h=None):
    p = ASSETS / name
    if p.exists():
        s.shapes.add_picture(str(p), l, t, w, h) if h else s.shapes.add_picture(str(p), l, t, w)

# ── Table helper ──
def add_table(s, left, top, width, height, data, col_widths=None, header_color=NAVY):
    rows, cols = len(data), len(data[0])
    tbl_shape = s.shapes.add_table(rows, cols, left, top, width, height)
    tbl = tbl_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            tbl.columns[i].width = w

    for r, row_data in enumerate(data):
        for c, cell_text in enumerate(row_data):
            cell = tbl.cell(r, c)
            cell.text = str(cell_text)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(9)
                p.font.name = "Calibri"
                if r == 0:
                    p.font.bold = True
                    p.font.color.rgb = WHITE
                else:
                    p.font.color.rgb = SL8
            # Header row coloring
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if r % 2 == 1 else RGBColor(0xF1, 0xF5, 0xF9)
    return tbl_shape


# ═══════════════ BUILD ═══════════════
def main():
    prs = Presentation()
    prs.slide_width = SW; prs.slide_height = SH
    blank = prs.slide_layouts[6]

    # ═══ 1. TITLE SLIDE ═══
    s = prs.slides.add_slide(blank); bg(s, NAVY_DK); add_tr(s, "fade")
    band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(1.8), SW, Inches(3.2))
    band.fill.solid(); band.fill.fore_color.rgb = NAVY; band.line.fill.background()

    tx(s, Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.4),
       "BE Project Presentation  |  Department of Computer Engineering", sz=12, c=SL4, al=PP_ALIGN.CENTER)

    tx(s, Inches(0.8), Inches(2.0), Inches(11.7), Inches(1.0),
       "Genome-Doc: Hallucination-Free Document\nRestoration via Symbolic Genome Inference",
       sz=32, c=WHITE, b=True, al=PP_ALIGN.CENTER, fn="Calibri Light")

    tx(s, Inches(0.8), Inches(3.4), Inches(11.7), Inches(0.4),
       "A Novel Structured Intermediate Representation for Neural Document Restoration",
       sz=14, c=SL2, al=PP_ALIGN.CENTER)

    tx(s, Inches(1.5), Inches(5.3), Inches(10.3), Inches(1.2),
       "Group Members:\n"
       "    1. [Student Name 1]    2. [Student Name 2]\n"
       "    3. [Student Name 3]    4. [Student Name 4]\n\n"
       "Guide:  Prof. [Guide Name]",
       sz=14, c=SL4, al=PP_ALIGN.CENTER)

    snum(s, 1)

    # ═══ 2. CONTENTS ═══
    s = prs.slides.add_slide(blank)
    header(s, "OVERVIEW", "Contents", n=2, tr="push")

    contents = [
        "Broad Area of Research", "Introduction", "Motivation",
        "Literature Review", "Existing System", "Research Gaps",
        "Proposed System & Architecture", "Objectives",
        "Algorithms Used", "Software & Hardware",
        "Datasets Used", "Implementation Screenshots",
        "Future Scope", "Conclusion", "References"
    ]
    for i, item in enumerate(contents):
        col = 0 if i < 8 else 1
        row = i if i < 8 else i - 8
        x = Inches(1.5 + col * 5.5)
        y = Inches(1.8 + row * 0.6)
        num_color = TEAL if col == 0 else INDIGO
        tx(s, x, y, Inches(0.5), Inches(0.35),
           f"{i+1:02d}", sz=13, c=num_color, b=True, fn="Consolas")
        tx(s, x + Inches(0.5), y, Inches(4.5), Inches(0.35),
           item, sz=13, c=SL8)

    # ═══ 3. BROAD AREA OF RESEARCH ═══
    s = prs.slides.add_slide(blank)
    header(s, "RESEARCH AREA", "Broad Area of Research", n=3, tr="wipe")

    ml(s, Inches(0.8), Inches(1.7), Inches(11), Inches(1.0), [
        ("This project falls at the intersection of:", SL6, False, 14),
    ])

    areas = [
        ("Document Image Analysis (DIA)", "Automated understanding of scanned/photographed documents — layout detection, text recognition, structural parsing."),
        ("Neural Image Restoration", "Using deep learning to recover clean images from degraded inputs — denoising, deblurring, super-resolution, inpainting."),
        ("Generative AI for Documents", "Applying diffusion models, GANs, and transformer-based generation to produce high-fidelity document images."),
        ("Structured Prediction & Neuro-Symbolic AI", "Combining neural perception (image understanding) with symbolic reasoning (structured JSON output) for verifiable results."),
    ]
    for i, (title, desc) in enumerate(areas):
        y = Inches(2.5 + i * 1.1)
        crd(s, Inches(0.8), y, Inches(11.5), Inches(0.9), bc=TEAL)
        topbar(s, Inches(0.8), y, Inches(11.5), TEAL)
        tx(s, Inches(1.1), y + Inches(0.1), Inches(11), Inches(0.3),
           title, sz=14, c=SL8, b=True)
        tx(s, Inches(1.1), y + Inches(0.45), Inches(11), Inches(0.35),
           desc, sz=11, c=SL6)

    # ═══ 4. INTRODUCTION ═══
    s = prs.slides.add_slide(blank)
    header(s, "INTRODUCTION", "Introduction", n=4, tr="cover")

    ml(s, Inches(0.8), Inches(1.7), Inches(7), Inches(5), [
        ("Document restoration is the task of recovering a clean, readable document from a degraded input — addressing noise, stains, fading, tears, and other damage.", SL6, False, 13),
        ("", SL6, False, 6),
        ("Genome-Doc proposes a paradigm shift: instead of pixel-to-pixel restoration, we extract a structured symbolic representation (the Document Genome) and then re-render a clean document from that specification.", SL6, False, 13),
        ("", SL6, False, 6),
        ("Key Insight:", NAVY, True, 14),
        ("By making the intermediate representation explicit and verifiable, we can guarantee that every character in the output was deliberately inferred — not hallucinated by a diffusion model.", SL6, False, 13),
        ("", SL6, False, 6),
        ("The system comprises three modules:", NAVY, True, 13),
        ("SIR — Style & Identity Refiner (captures visual DNA)", SL6, False, 12),
        ("DGI — Document Genome Inferrer (extracts structured content)", SL6, False, 12),
        ("NRE — Neural Re-Rendering Engine (generates clean output)", SL6, False, 12),
    ])

    add_img(s, "genome_hero.png", Inches(8.5), Inches(1.7), Inches(4.3), Inches(4.3))

    # ═══ 5. MOTIVATION ═══
    s = prs.slides.add_slide(blank)
    header(s, "MOTIVATION", "Motivation", n=5, tr="reveal")

    motivations = [
        ("Text Hallucination in Existing Methods",
         "Current diffusion-based restorers (DE-GAN, DocDiff, DocRes) treat text as pixels — they generate visually plausible but factually wrong characters. A single wrong digit in a legal document changes its meaning entirely."),
        ("No Structural Guarantees",
         "End-to-end models have no mechanism to verify that restored text matches the original. There is no intermediate check between input and output."),
        ("Style Destruction",
         "Existing methods replace original fonts, ink weight, and paper texture with generic defaults, destroying the document's visual identity and provenance."),
        ("Archival & Legal Requirements",
         "Historical archives, legal records, and government documents demand verifiable restoration — not statistical guesses. The output must be traceable to explicit visual evidence."),
    ]
    for i, (title, desc) in enumerate(motivations):
        y = Inches(1.7 + i * 1.3)
        crd(s, Inches(0.8), y, Inches(11.5), Inches(1.1))
        topbar(s, Inches(0.8), y, Inches(11.5), ROSE if i == 0 else [AMBER6, INDIGO, TEAL][i-1])
        tx(s, Inches(1.1), y + Inches(0.12), Inches(11), Inches(0.3),
           f"{i+1}. {title}", sz=14, c=SL8, b=True)
        tx(s, Inches(1.1), y + Inches(0.5), Inches(11), Inches(0.45),
           desc, sz=11, c=SL6)

    # ═══ 6. LITERATURE REVIEW TABLE ═══
    s = prs.slides.add_slide(blank)
    header(s, "LITERATURE REVIEW", "Literature Review", n=6, tr="push")

    lit_data = [
        ["Ref.", "Paper Title", "Authors & Year", "Techniques", "Key Highlights", "Dataset", "Limitations"],
        ["[1]", "DE-GAN", "Souibgui et al., 2020", "GAN-based\nimage-to-image", "First GAN for\ndocument enhancement", "DIBCO,\nH-DIBCO", "Hallucinates text\nin heavy degradation"],
        ["[2]", "DocDiff", "Yang et al., 2023", "Diffusion model\nwith residual pred.", "Coarse-to-fine\ndenoising", "DocUNet,\nDIBCO", "No structural\nguarantees on text"],
        ["[3]", "DocRes", "Zhang et al., 2024", "Multi-task\nrestoration", "Unified model for\nmultiple degradations", "DIBCO,\ncustom", "Generic style,\nloses identity"],
        ["[4]", "Donut", "Kim et al., 2022", "OCR-free VL\ntransformer", "No OCR needed,\nend-to-end", "CORD,\nRVL-CDIP", "No layout\nregression head"],
        ["[5]", "ControlNet", "Zhang et al., 2023", "Conditional\ndiffusion control", "Spatial conditioning\nfor SD", "LAION,\ncustom", "Struggles with\ncrisp text rendering"],
    ]
    add_table(s, Inches(0.4), Inches(1.7), Inches(12.5), Inches(5.2),
              lit_data,
              col_widths=[Inches(0.5), Inches(1.3), Inches(1.8), Inches(1.7), Inches(2.2), Inches(1.5), Inches(2.0)])

    # ═══ 7. EXISTING SYSTEM ═══
    s = prs.slides.add_slide(blank)
    header(s, "EXISTING SYSTEM", "Existing System", n=7, tr="wipe")

    tx(s, Inches(0.8), Inches(1.7), Inches(11), Inches(0.5),
       "Current document restoration follows a direct image-to-image translation paradigm:",
       sz=14, c=SL6)

    crd(s, Inches(0.8), Inches(2.4), Inches(11.5), Inches(1.8))
    tx(s, Inches(1.2), Inches(2.6), Inches(10.5), Inches(0.3),
       "Image-to-Image Pipeline (DE-GAN, DocDiff, DocRes)", sz=16, c=NAVY, b=True)

    ml(s, Inches(1.2), Inches(3.1), Inches(10.5), Inches(1.0), [
        ("Degraded Image → Neural Network (GAN/Diffusion) → Clean Image", SL8, False, 13),
        ("No intermediate representation — the network directly maps pixels to pixels", SL6, False, 12),
        ("Text content is treated as visual texture, not semantic information", SL6, False, 12),
    ])

    tx(s, Inches(0.8), Inches(4.5), Inches(11), Inches(0.3),
       "Key Limitations of Existing Systems:", sz=15, c=NAVY, b=True)

    lims = [
        "No mechanism to verify whether restored text is correct or hallucinated",
        "Single wrong character can change meaning (e.g., '1000' → '1900' in a historical ledger)",
        "Original document style (font, paper, ink) replaced by generic learned distribution",
        "No structured output — cannot extract text, layout, or metadata from restored image",
        "Performance degrades severely on unseen degradation types",
    ]
    bullets(s, Inches(1.0), Inches(5.0), Inches(11), lims, sz=12, c=SL6, gap=0.38)

    # ═══ 8. RESEARCH GAPS ═══
    s = prs.slides.add_slide(blank)
    header(s, "RESEARCH GAPS", "Research Gaps", n=8, tr="split")

    gaps = [
        ("Gap 1: No Anti-Hallucination Mechanism",
         "Existing methods lack any verification step between input and output. The restored text is a statistical guess, not a verified reconstruction.",
         "Genome-Doc introduces an explicit, inspectable Document Genome as intermediate representation."),
        ("Gap 2: Style-Agnostic Restoration",
         "Current models learn a single 'average' document appearance. The unique visual identity of each document is lost during restoration.",
         "SIR module extracts document-specific style embeddings to condition the renderer."),
        ("Gap 3: No Structural Understanding",
         "Pixel-level models cannot distinguish headings from body text, tables from figures. Layout structure is not preserved.",
         "DGI outputs a structured JSON with semantic types (heading, paragraph, table) and bounding boxes."),
        ("Gap 4: Non-Verifiable Output",
         "There is no way to audit what the model 'decided' — the restored image is an opaque neural output.",
         "The Document Genome is human-readable JSON that can be inspected and corrected before rendering."),
    ]
    for i, (title, gap_desc, our_sol) in enumerate(gaps):
        y = Inches(1.7 + i * 1.35)
        crd(s, Inches(0.8), y, Inches(5.4), Inches(1.15), bc=ROSE)
        topbar(s, Inches(0.8), y, Inches(5.4), ROSE)
        tx(s, Inches(1.1), y + Inches(0.08), Inches(4.8), Inches(0.3),
           title, sz=12, c=SL8, b=True)
        tx(s, Inches(1.1), y + Inches(0.4), Inches(4.8), Inches(0.6),
           gap_desc, sz=10, c=SL6)

        crd(s, Inches(6.5), y, Inches(5.8), Inches(1.15), bc=GREEN7)
        topbar(s, Inches(6.5), y, Inches(5.8), GREEN7)
        tx(s, Inches(6.8), y + Inches(0.08), Inches(5.2), Inches(0.3),
           "Our Solution", sz=12, c=GREEN7, b=True)
        tx(s, Inches(6.8), y + Inches(0.4), Inches(5.2), Inches(0.6),
           our_sol, sz=10, c=SL6)

    # ═══ 9. PROPOSED SYSTEM & ARCHITECTURE ═══
    s = prs.slides.add_slide(blank)
    header(s, "PROPOSED SYSTEM", "Proposed System — Methodology & Architecture", n=9, tr="cover")

    tx(s, Inches(0.8), Inches(1.7), Inches(11), Inches(0.5),
       "Three-module pipeline: DGI extracts the Document Genome, SIR captures style, NRE re-renders the clean document.",
       sz=13, c=SL6)

    # Pipeline boxes
    pipeline = [
        ("Degraded\nDocument", "Input", SL6),
        ("DGI", "Genome\nExtraction", INDIGO),
        ("SIR", "Style\nCapture", TEAL),
        ("NRE", "Neural\nRe-Rendering", ROSE),
        ("Clean\nDocument", "Output", GREEN7),
    ]
    for i, (name, desc, col) in enumerate(pipeline):
        x = Inches(0.5 + i * 2.6)
        crd(s, x, Inches(2.5), Inches(2.2), Inches(2.0), bc=col)
        topbar(s, x, Inches(2.5), Inches(2.2), col)
        tx(s, x + Inches(0.1), Inches(2.7), Inches(2.0), Inches(0.7),
           name, sz=16, c=col, b=True, al=PP_ALIGN.CENTER)
        tx(s, x + Inches(0.1), Inches(3.5), Inches(2.0), Inches(0.5),
           desc, sz=11, c=SL4, al=PP_ALIGN.CENTER)
        if i < 4:
            tx(s, x + Inches(2.2), Inches(3.1), Inches(0.4), Inches(0.4),
               "→", sz=20, c=SL4, al=PP_ALIGN.CENTER)

    # Module details
    modules = [
        ("SIR (Style & Identity Refiner)", TEAL,
         ["ResNet-50 backbone + attention pooling", "InfoNCE contrastive loss",
          "Outputs 512-dim style embedding", "Status: COMPLETE"]),
        ("DGI (Document Genome Inferrer)", INDIGO,
         ["Donut vision-language backbone", "Template-guided JSON decoding",
          "LayoutHead for bbox regression", "LoRA fine-tuning (2.3% params)"]),
        ("NRE (Neural Re-Rendering Engine)", ROSE,
         ["ControlNet + Stable Diffusion 1.5", "Style injection via cross-attention",
          "Skeleton map as spatial condition", "DDIM 50-step inference"]),
    ]
    for i, (title, col, items) in enumerate(modules):
        x = Inches(0.8 + i * 4.1)
        crd(s, x, Inches(4.8), Inches(3.8), Inches(2.3), bc=col)
        topbar(s, x, Inches(4.8), Inches(3.8), col)
        tx(s, x + Inches(0.2), Inches(4.95), Inches(3.4), Inches(0.3),
           title, sz=12, c=col, b=True)
        for j, item in enumerate(items):
            c_item = GREEN7 if "COMPLETE" in item else SL6
            tx(s, x + Inches(0.2), Inches(5.35 + j * 0.38), Inches(3.4), Inches(0.3),
               f"•  {item}", sz=10, c=c_item)

    # ═══ 10. OBJECTIVES ═══
    s = prs.slides.add_slide(blank)
    header(s, "OBJECTIVES", "Objectives", n=10, tr="fade")

    objectives = [
        "To develop a hallucination-free document restoration system that guarantees every character in the output is explicitly inferred from the input image.",
        "To design and implement the Document Genome — a structured JSON intermediate representation capturing text content, layout bounding boxes, and semantic element types.",
        "To build a Style & Identity Refiner (SIR) that captures document-specific visual attributes (font, ink, paper texture) as a compact 512-dimensional embedding.",
        "To create a Document Genome Inferrer (DGI) using a Donut vision-language backbone with template-guided decoding and LoRA fine-tuning.",
        "To implement a Neural Re-Rendering Engine (NRE) using ControlNet-conditioned Stable Diffusion with style embedding injection for photorealistic output.",
        "To evaluate the system against standard benchmarks (DIBCO) using PSNR, SSIM, CER, and a novel hallucination rate metric.",
    ]
    for i, obj in enumerate(objectives):
        y = Inches(1.8 + i * 0.85)
        # Number badge
        sq = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(0.8), y, Inches(0.45), Inches(0.45))
        sq.fill.solid(); sq.fill.fore_color.rgb = NAVY; sq.line.fill.background()
        p = sq.text_frame.paragraphs[0]; p.text = str(i+1)
        p.font.size = Pt(14); p.font.color.rgb = WHITE; p.font.bold = True
        p.font.name = "Consolas"; p.alignment = PP_ALIGN.CENTER
        sq.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        # Text
        tx(s, Inches(1.45), y + Inches(0.02), Inches(11), Inches(0.5),
           obj, sz=12, c=SL6)

    # ═══ 11. ALGORITHMS USED ═══
    s = prs.slides.add_slide(blank)
    header(s, "ALGORITHMS", "Algorithms Used in Proposed System", n=11, tr="push")

    algo_data = [
        ["Algorithm", "Module", "Purpose", "Details"],
        ["InfoNCE\nContrastive Loss", "SIR", "Style embedding\ntraining", "Pulls same-document patches together,\npushes different docs apart in latent space"],
        ["ResNet-50\n+ Attention Pooling", "SIR", "Feature extraction\n& aggregation", "ImageNet-pretrained backbone with attention-\nweighted mean pooling over K patches"],
        ["Donut (VL\nTransformer)", "DGI", "Genome\nextraction", "Swin encoder → BART decoder,\nOCR-free vision-language architecture"],
        ["LoRA (Low-Rank\nAdaptation)", "DGI", "Parameter-\nefficient fine-tuning", "Rank-16 adapters on decoder attention;\nonly 2.3% of params trained"],
        ["ControlNet +\nDDIM", "NRE", "Conditional\nimage generation", "Skeleton map conditions SD 1.5;\n50-step DDIM denoising at inference"],
        ["GIoU Loss", "DGI", "Bounding box\nregression", "Generalized IoU loss + L1 for\nlayout coordinate prediction"],
    ]
    add_table(s, Inches(0.4), Inches(1.7), Inches(12.5), Inches(5.0), algo_data,
              col_widths=[Inches(2.0), Inches(1.0), Inches(2.0), Inches(5.5)])

    # ═══ 12. SOFTWARE/HARDWARE ═══
    s = prs.slides.add_slide(blank)
    header(s, "INFRASTRUCTURE", "Software & Hardware", n=12, tr="wipe")

    tx(s, Inches(0.8), Inches(1.7), Inches(5), Inches(0.3),
       "Software Stack", sz=18, c=NAVY, b=True)

    sw_items = [
        ["Component", "Technology", "Version"],
        ["Language", "Python", "3.10+"],
        ["Deep Learning", "PyTorch", "2.1+"],
        ["Vision Backbone", "torchvision (ResNet-50)", "0.16+"],
        ["VL Model", "Donut (HuggingFace)", "1.0"],
        ["Diffusion", "Stable Diffusion 1.5", "—"],
        ["ControlNet", "diffusers library", "0.25+"],
        ["Image Processing", "Pillow, OpenCV", "—"],
        ["Data Validation", "Pydantic", "2.0+"],
        ["Training", "FP16 AMP, AdamW", "—"],
    ]
    add_table(s, Inches(0.8), Inches(2.1), Inches(5.5), Inches(4.5), sw_items,
              col_widths=[Inches(1.8), Inches(2.2), Inches(1.0)])

    tx(s, Inches(7.0), Inches(1.7), Inches(5), Inches(0.3),
       "Hardware Requirements", sz=18, c=NAVY, b=True)

    hw_items = [
        ["Component", "Specification"],
        ["GPU (Training)", "NVIDIA RTX 6000 Pro\n(96 GB VRAM)"],
        ["GPU (Testing)", "NVIDIA T4 (16 GB)\non Google Colab"],
        ["GPU (Inference)", "Any CUDA GPU\n(≥8 GB VRAM)"],
        ["RAM", "32 GB minimum"],
        ["Storage", "50 GB for datasets\n+ checkpoints"],
        ["OS", "Linux / Windows 10+"],
    ]
    add_table(s, Inches(7.0), Inches(2.1), Inches(5.5), Inches(4.5), hw_items,
              col_widths=[Inches(2.0), Inches(3.0)])

    # ═══ 13. DATASETS ═══
    s = prs.slides.add_slide(blank)
    header(s, "DATASETS", "Datasets Used for Implementation", n=13, tr="cover")

    ds_data = [
        ["Dataset", "Size", "Contents", "Used For"],
        ["Custom Synthetic\n(Generated)", "30,000\npairs", "Clean docs + degraded versions\n+ ground-truth Genome JSONs", "Primary training\ndata for all modules"],
        ["DIBCO\n(2009–2019)", "~1,000\nimages", "Standard binarization benchmark;\nhistorical document images", "Evaluation against\nprior work"],
        ["Google Fonts", "1,500+\nfamilies", "Open-source font collection\nfor realistic rendering", "Synthetic data\ngeneration"],
        ["DocLayNet\n(IBM)", "80,000\npages", "Layout annotations for\n6 document categories", "Layout distribution\nstatistics"],
    ]
    add_table(s, Inches(0.8), Inches(1.7), Inches(11.5), Inches(3.5), ds_data,
              col_widths=[Inches(2.0), Inches(1.2), Inches(4.2), Inches(3.0)])

    tx(s, Inches(0.8), Inches(5.5), Inches(11), Inches(0.3),
       "Synthetic Data Generation Pipeline:", sz=14, c=NAVY, b=True)

    ml(s, Inches(0.8), Inches(5.9), Inches(11), Inches(1.0), [
        ("•  Random layout sampling from DocLayNet statistics (headers, paragraphs, tables)", SL6, False, 11),
        ("•  Text rendering with 1,500+ Google Fonts at varying sizes, weights, and colors", SL6, False, 11),
        ("•  8 degradation types: Gaussian noise, blur, stains, yellowing, JPEG compression, warping, ink bleed, creases", SL6, False, 11),
    ])

    # ═══ 14. IMPLEMENTATION SCREENSHOTS ═══
    s = prs.slides.add_slide(blank)
    header(s, "IMPLEMENTATION", "Implementation Screenshots", n=14, tr="reveal")

    tx(s, Inches(0.8), Inches(1.7), Inches(11), Inches(0.5),
       "Screenshots will be added after model training is complete. Current status:",
       sz=14, c=SL6)

    statuses = [
        ("SIR Module — Training", "COMPLETE", GREEN7, GREEN1),
        ("DGI Module — Training", "IN PROGRESS", AMBER6, AMBER1),
        ("NRE Module — Training", "PENDING", SL4, SL2),
        ("End-to-End Inference", "PENDING", SL4, SL2),
    ]
    for i, (label, status, fg, bgc) in enumerate(statuses):
        y = Inches(2.5 + i * 1.0)
        crd(s, Inches(1.5), y, Inches(10.3), Inches(0.8))
        tx(s, Inches(2.0), y + Inches(0.2), Inches(6), Inches(0.4),
           label, sz=15, c=SL8, b=True)
        tag(s, Inches(9.5), y + Inches(0.22), status, fg, bgc, Inches(1.8))

    tx(s, Inches(1.5), Inches(6.2), Inches(10), Inches(0.4),
       "Placeholder: training logs, sample restorations, and metric visualizations will be inserted post-training.",
       sz=11, c=SL4, al=PP_ALIGN.CENTER)

    # ═══ 15. FUTURE SCOPE ═══
    s = prs.slides.add_slide(blank)
    header(s, "FUTURE SCOPE", "Future Scope", n=15, tr="fade")

    futures = [
        ("Multi-Language Support",
         "Extend the DGI tokenizer and Genome schema to handle non-Latin scripts (Devanagari, Arabic, CJK) for global applicability."),
        ("Real-Time Processing",
         "Optimize the inference pipeline with model distillation and quantization (INT8/FP16) for real-time document scanning applications."),
        ("Handwritten Document Support",
         "Adapt the SIR and DGI modules for handwritten text recognition and style capture, extending beyond printed documents."),
        ("Interactive Genome Editing",
         "Build a GUI that displays the inferred Document Genome for human review and correction before re-rendering — enabling human-in-the-loop restoration."),
        ("End-to-End Joint Training",
         "Implement differentiable genome generation using Gumbel-Softmax relaxation to enable true end-to-end backpropagation across all three modules."),
        ("Mobile Deployment",
         "Develop lightweight model variants using MobileNet/EfficientNet backbones for on-device document scanning and restoration."),
    ]
    for i, (title, desc) in enumerate(futures):
        col = 0 if i < 3 else 1
        row = i if i < 3 else i - 3
        x = Inches(0.8 + col * 6.2)
        y = Inches(1.8 + row * 1.7)
        crd(s, x, y, Inches(5.8), Inches(1.5))
        topbar(s, x, y, Inches(5.8), [TEAL, INDIGO, ROSE, NAVY, AMBER6, GREEN7][i])
        tx(s, x + Inches(0.3), y + Inches(0.15), Inches(5.2), Inches(0.3),
           title, sz=14, c=SL8, b=True)
        tx(s, x + Inches(0.3), y + Inches(0.55), Inches(5.2), Inches(0.7),
           desc, sz=11, c=SL6)

    # ═══ 16. CONCLUSION ═══
    s = prs.slides.add_slide(blank)
    header(s, "CONCLUSION", "Conclusion", n=16, tr="wipe")

    conclusions = [
        "Genome-Doc introduces a novel paradigm for document restoration: symbolic genome inference followed by conditional neural re-rendering, eliminating hallucination by construction.",
        "The Document Genome provides a human-readable, structured intermediate representation that can be independently verified before the final image is generated.",
        "The SIR module (complete) successfully captures document-specific style embeddings using contrastive learning, enabling style-faithful restoration.",
        "The three-module architecture (DGI → SIR → NRE) separates content understanding, style extraction, and image generation into independently trainable and verifiable components.",
        "Expected metrics: PSNR ≥ 28 dB, SSIM ≥ 0.92, CER < 5%, and near-zero hallucination rate on standard benchmarks.",
        "This approach opens new directions in verifiable AI — where the reasoning process (the Genome) is as important as the final output.",
    ]
    for i, conc in enumerate(conclusions):
        y = Inches(1.7 + i * 0.88)
        sq = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(0.8), y + Inches(0.02), Inches(0.35), Inches(0.35))
        sq.fill.solid(); sq.fill.fore_color.rgb = TEAL; sq.line.fill.background()
        p = sq.text_frame.paragraphs[0]; p.text = str(i+1)
        p.font.size = Pt(12); p.font.color.rgb = WHITE; p.font.bold = True
        p.font.name = "Consolas"; p.alignment = PP_ALIGN.CENTER
        sq.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        tx(s, Inches(1.35), y, Inches(11), Inches(0.7),
           conc, sz=12, c=SL6)

    # ═══ 17. REFERENCES ═══
    s = prs.slides.add_slide(blank)
    header(s, "REFERENCES", "References", n=17, tr="push")

    refs = [
        "[1]  Souibgui, M. A., & Kessentini, Y. (2020). DE-GAN: A Conditional Generative Adversarial Network for Document Enhancement. IEEE Trans. PAMI.",
        "[2]  Yang, Z., et al. (2023). DocDiff: Document Enhancement via Residual Denoising Diffusion Probabilistic Models. ACM Multimedia.",
        "[3]  Zhang, J., et al. (2024). DocRes: A Generalist Model Toward Unifying Document Image Restoration Tasks. CVPR.",
        "[4]  Kim, G., et al. (2022). OCR-free Document Understanding Transformer (Donut). ECCV.",
        "[5]  Zhang, L., et al. (2023). Adding Conditional Control to Text-to-Image Diffusion Models (ControlNet). ICCV.",
        "[6]  Hu, E. J., et al. (2022). LoRA: Low-Rank Adaptation of Large Language Models. ICLR.",
        "[7]  Oord, A. v. d., Li, Y., & Vinyals, O. (2018). Representation Learning with Contrastive Predictive Coding (InfoNCE). arXiv.",
        "[8]  Pratikakis, I., et al. (2019). ICDAR 2019 Competition on Document Image Binarization (DIBCO). ICDAR.",
        "[9]  Rombach, R., et al. (2022). High-Resolution Image Synthesis with Latent Diffusion Models (Stable Diffusion). CVPR.",
        "[10] Pfister, A., et al. (2022). DocLayNet: A Large Human-Annotated Dataset for Document Layout Analysis. KDD.",
    ]
    for i, ref in enumerate(refs):
        col = 0 if i < 5 else 1
        row = i if i < 5 else i - 5
        x = Inches(0.8 + col * 6.3)
        y = Inches(1.8 + row * 1.0)
        tx(s, x, y, Inches(5.8), Inches(0.8),
           ref, sz=10, c=SL6)

    # ── Save ──
    out = DIR / "GenomeDoc_Presentation.pptx"
    prs.save(str(out))
    print(f"✅ Saved → {out}")


if __name__ == "__main__":
    main()
